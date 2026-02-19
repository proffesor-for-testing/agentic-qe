#!/usr/bin/env python3
"""
Create a PPTX with one full-screen auto-playing video per slide.
Videos auto-advance so in presentation mode it looks like one continuous video.
"""

import os
import glob
import struct
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn, nsmap
from lxml import etree

SLIDE_DIR = "/workspaces/agentic-qe/Agentic QCSD/03 Reference Docs/slide-videos"
OUTPUT = "/workspaces/agentic-qe/Agentic QCSD/03 Reference Docs/Agentic-QCSD-Teaser.pptx"

# 16:9 widescreen
SLIDE_W = Emu(12192000)  # 13.333 in
SLIDE_H = Emu(6858000)   # 7.5 in

# Slide titles for reference (not displayed)
SLIDE_TITLES = [
    "Opening Hook",
    "Problem Stats",
    "Faster Horses",
    "Introducing QCSD",
    "The Framework",
    "Award Credibility",
    "Five Swarms",
    "Feedback Loops",
    "Architecture Engine",
    "The Outcome",
    "Claude Flow Plugin",
    "Credits",
    "CTA Closing",
]

# Durations in ms (must match recording durations for auto-advance)
SLIDE_DURATIONS_MS = [
    8000, 12000, 6000, 8000, 15000, 10000, 14000,
    12000, 14000, 10000, 12000, 10000, 10000,
]


def get_video_files():
    """Get sorted MP4 files."""
    files = sorted(glob.glob(os.path.join(SLIDE_DIR, "*.mp4")))
    return files


def set_black_background(slide):
    """Set slide background to solid black via direct XML."""
    from pptx.dml.color import RGBColor
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)


def add_video_to_slide(prs, slide, video_path, duration_ms):
    """Add a full-screen, auto-playing video to a slide."""
    # Add video as a movie shape
    movie = slide.shapes.add_movie(
        video_path,
        left=0, top=0,
        width=SLIDE_W, height=SLIDE_H,
        mime_type='video/mp4',
    )

    # Get the pic element (the shape XML)
    pic = movie._element

    # -- Make video auto-play on slide load --
    # We need to add timing info to the slide XML
    slide_elem = slide._element

    # Build the timing tree for auto-play
    # <p:timing>
    #   <p:tnLst>
    #     <p:par>
    #       <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
    #         <p:childTnLst>
    #           <p:seq concurrent="1" nextAc="seek">
    #             <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
    #               <p:childTnLst>
    #                 <p:par>
    #                   <p:cTn id="3" fill="hold">
    #                     <p:stCondLst><p:cond delay="0"/></p:stCondLst>
    #                     <p:childTnLst>
    #                       <p:par>
    #                         <p:cTn id="4" fill="hold">
    #                           <p:stCondLst><p:cond delay="0"/></p:stCondLst>
    #                           <p:childTnLst>
    #                             <p:cmd type="call" cmd="playFrom(0)">
    #                               ... target video shape ...

    # Get the shape id for targeting
    sp_id = pic.find(qn('p:nvPicPr')).find(qn('p:cNvPr')).get('id')

    timing_xml = f'''<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                               xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="3" fill="hold">
                        <p:stCondLst>
                          <p:cond delay="0"/>
                        </p:stCondLst>
                        <p:childTnLst>
                          <p:par>
                            <p:cTn id="4" fill="hold">
                              <p:stCondLst>
                                <p:cond delay="0"/>
                              </p:stCondLst>
                              <p:childTnLst>
                                <p:cmd type="call" cmd="playFrom(0)">
                                  <p:cBhvr>
                                    <p:cTn id="5" dur="{duration_ms}" fill="hold"/>
                                    <p:tgtEl>
                                      <p:spTgt spid="{sp_id}"/>
                                    </p:tgtEl>
                                  </p:cBhvr>
                                </p:cmd>
                              </p:childTnLst>
                            </p:cTn>
                          </p:par>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst>
                  <p:cond evt="onPrev" delay="0">
                    <p:tgtEl><p:sldTgt/></p:tgtEl>
                  </p:cond>
                </p:prevCondLst>
                <p:nextCondLst>
                  <p:cond evt="onNext" delay="0">
                    <p:tgtEl><p:sldTgt/></p:tgtEl>
                  </p:cond>
                </p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>'''

    timing_elem = etree.fromstring(timing_xml)

    # Remove any existing timing
    existing = slide_elem.find(qn('p:timing'))
    if existing is not None:
        slide_elem.remove(existing)

    slide_elem.append(timing_elem)

    # -- Auto-advance slide after video duration --
    # Set advTm on the slide transition (milliseconds)
    transition_xml = f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="fast" advClick="0" advTm="{duration_ms}"/>'
    transition_elem = etree.fromstring(transition_xml)

    existing_trans = slide_elem.find(qn('p:transition'))
    if existing_trans is not None:
        slide_elem.remove(existing_trans)

    # Insert transition before timing
    timing_idx = list(slide_elem).index(slide_elem.find(qn('p:timing')))
    slide_elem.insert(timing_idx, transition_elem)


def main():
    video_files = get_video_files()
    print(f"Found {len(video_files)} video files")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    for i, vf in enumerate(video_files):
        title = SLIDE_TITLES[i] if i < len(SLIDE_TITLES) else f"Slide {i+1}"
        duration = SLIDE_DURATIONS_MS[i] if i < len(SLIDE_DURATIONS_MS) else 10000
        print(f"  Adding slide {i+1}: {title} ({os.path.basename(vf)}, {duration/1000}s)")

        slide = prs.slides.add_slide(blank_layout)
        set_black_background(slide)
        add_video_to_slide(prs, slide, vf, duration)

    prs.save(OUTPUT)
    size_mb = os.path.getsize(OUTPUT) / (1024 * 1024)
    print(f"\nSaved: {OUTPUT}")
    print(f"Size: {size_mb:.1f} MB")
    print(f"Slides: {len(prs.slides)}")
    print("Auto-play + auto-advance enabled on all slides.")


if __name__ == "__main__":
    main()
